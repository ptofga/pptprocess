# -*- coding: utf-8 -*-
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import sys
import threading
import queue
from pptx import Presentation
import pandas as pd

class AsyncRedirectText:
    def __init__(self, text_widget):
        self.text_widget = text_widget
        self.queue = queue.Queue()
        self.running = True

    def write(self, message):
        self.queue.put(message)

    def flush(self): 
        pass

    def start_polling(self):
        def check_queue():
            while not self.queue.empty():
                message = self.queue.get()
                self.text_widget.configure(state='normal')
                self.text_widget.insert(tk.END, message)
                self.text_widget.see(tk.END)
                self.text_widget.configure(state='disabled')
            if self.running:
                self.text_widget.after(100, check_queue)
        check_queue()

    def stop(self):
        self.running = False

class ModernPPTXApp:
    def __init__(self, root):
        self.root = root
        self.root.title("SPR Curve 数据分析 V3")
        self.root.geometry("600x400")
        self.style = ttk.Style()
        self.configure_styles()
        
        # 初始化变量
        self.pptx_path = tk.StringVar()
        self.running = False
        
        # 创建 UI 组件
        self.create_widgets()
        
        # 异步输出重定向
        self.async_redirect = AsyncRedirectText(self.output_text)
        self.async_redirect.start_polling()
        sys.stdout = self.async_redirect
        self.sample_count = tk.IntVar(value=48)  # 默认采样数
        self.sample_entry.config(textvariable=self.sample_count)

    def configure_styles(self):
        """配置现代UI样式"""
        self.style.theme_use('clam')
        self.style.configure('TButton', font=('微软雅黑', 10), padding=6)
        self.style.configure('TLabel', font=('微软雅黑', 9))
        self.style.configure('Header.TFrame', background='#f0f0f0')
        self.style.map('TButton',
            foreground=[('active', '#ffffff'), ('!active', '#333333')],
            background=[('active', '#0078d7'), ('!active', '#f0f0f0')]
        )
    # ... [保持其他方法不变] ...
    def create_widgets(self):
        # 主容器
        main_frame = ttk.Frame(self.root)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        # 头部控制面板
        header_frame = ttk.Frame(main_frame, style='Header.TFrame')
        header_frame.pack(fill=tk.X, pady=(0,10))

        ttk.Button(header_frame, text="📂 选择PPT文件", 
                 command=self.select_pptx).pack(side=tk.LEFT, padx=5)
        ttk.Button(header_frame, text="⚙️ 开始分析", 
                 command=self.process_pptx).pack(side=tk.LEFT, padx=5)
        ttk.Button(header_frame, text="🧹 清空输出", 
                 command=self.clear_output).pack(side=tk.RIGHT, padx=5)
        
        ttk.Button(header_frame, text="软件说明", 
                 command=self.software_discription).pack(side=tk.RIGHT, padx=5)
        

        # 文件路径显示
        path_frame = ttk.Frame(main_frame)
        path_frame.pack(fill=tk.X, pady=5)

        # how many values sampled
        ttk.Label(path_frame, text="数据采样数:", width=10).pack(side=tk.LEFT, padx=(0,5))
        self.sample_entry = ttk.Entry(path_frame, width=10)
        self.sample_entry.pack(side=tk.LEFT)
        ttk.Label(path_frame, text="  ").pack(side=tk.LEFT)  # 添加间距
         # 在self.sample_entry.pack之后添加验证
        self.sample_entry.config(
            validate="key",
            validatecommand=(self.root.register(self.validate_number), '%P')
        )

        ttk.Label(path_frame, text="当前文件:").pack(side=tk.LEFT)
        self.path_label = ttk.Label(path_frame, textvariable=self.pptx_path, 
                                  foreground="#666666", width=60)
        self.path_label.pack(side=tk.LEFT, padx=5)

        # 输出区域
        output_frame = ttk.LabelFrame(main_frame, text="分析结果", padding=10)
        output_frame.pack(fill=tk.BOTH, expand=True)

        # 带滚动条的文本框
        text_scroll = ttk.Scrollbar(output_frame)
        text_scroll.pack(side=tk.RIGHT, fill=tk.Y)

        self.output_text = tk.Text(
            output_frame, 
            wrap=tk.WORD,
            yscrollcommand=text_scroll.set,
            font=('Consolas', 9),
            bg='#ffffff',
            padx=10,
            pady=10
        )
        self.output_text.pack(fill=tk.BOTH, expand=True)
        text_scroll.config(command=self.output_text.yview)

        # 状态栏
        self.status_bar = ttk.Label(
            main_frame, 
            text="就绪", 
            anchor=tk.W,
            foreground="#666666"
        )
        self.status_bar.pack(fill=tk.X, pady=(5,0))
    # 添加验证方法
    def validate_number(self, text):
        return text.isdigit() or text == ""
    
    def select_pptx(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("PowerPoint 文件", "*.pptx")],
            title="选择 PowerPoint 文件"
        )
        if file_path:
            self.pptx_path.set(file_path)
            self.status_bar.config(text=f"已选择文件: {file_path}")  
    def read_column_to_list(self,file_path, column_name=None, column_index=None, sheet_name=0):
        """
        读取 Excel 文件的某一列到列表
        :param file_path: Excel 文件路径
        :param column_name: 列名（如 "姓名"）
        :param column_index: 列索引（从 0 开始，如 2 表示第 3 列）
        :param sheet_name: 工作表名称或索引（默认为第一个工作表）
        :return: 列数据列表
        """
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        
        # 根据列名或列索引读取数据
        if column_name:
            column_data = df[column_name].tolist()
        elif column_index is not None:
            column_data = df.iloc[:, column_index].tolist()
        else:
            raise ValueError("必须指定 column_name 或 column_index")
        
        return column_data
 
    def retain_numbers_uniform(self,data,n):
        if not data:
            return []
        cleaned_list = [x for x in data if x!=' ' ]
        data_f = [float(x) for x in cleaned_list]

        m = len(data_f)
        if n <= 0 or m < n:
            raise ValueError("n必须大于0且不超过数据长度")
        if n == 1:
            return [data_f[0]]  # 或返回 [data[m//2]] 取中间值
        
        step = (m - 1) / (n - 1)
        indices = [int(round(i * step)) for i in range(n)]
        return [data_f[i] for i in indices]

    def retain_numbers_before_last_max(self,data,numbers):
        
        if not data:
            return []
        cleaned_list = [x for x in data if x!=' ' ]
        data_f = [float(x) for x in cleaned_list]
        max_value = max(data_f)
        # 找到所有最大值的索引
        max_indices = [i for i, x in enumerate(data_f) if x == max_value]
        if not max_indices:
            return []
        
        # 取最后一个最大值的索引
        max_index = max_indices[-1] 
        start_index = max(max_index - (numbers-1), 0)

        if(max_index<numbers):         
            return data_f[start_index : numbers]
        else:
            return data_f[start_index : max_index+1]

    def extract_chart_data(self,pptx_file):

        prs = Presentation(pptx_file)
        data = []  
        data_original =[]  
        columns = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    title=chart.chart_title.text_frame.text
                    no = title.split(';')[0]
                    print(no)
                    columns.append(no)
                    for series in chart.series:
                        data_original.append(series.values)
                        data.append(self.retain_numbers_before_last_max(series.values,int(self.sample_count.get())) ) 

        df = pd.DataFrame({columns[i]: data[i] for i in range(len(data))})
        df_orginal = pd.DataFrame({columns[i]: data_original[i] for i in range(len(data_original))})
        #df = pd.DataFrame(data)
        return df,df_orginal

    def calculate_mse(self,y_true, y_pred,distance):
        """
        计算两组数据的均方差（MSE）
        :param y_true: 实际值列表
        :param y_pred: 预测值列表
        :return: 均方差值
        """
        # 检查数据长度是否一致
        if len(y_true) != len(y_pred):
            raise ValueError("两组数据长度必须一致")
        
        # 计算平方差的和
        squared_errors = [(yt - yp+distance) ** 2 for yt, yp in zip(y_true, y_pred)]
        mse = sum(squared_errors) / len(y_true)
        return mse
    def process_pptx(self):
        if self.running:
            messagebox.showwarning("警告", "已有任务正在运行中")
            return
            
        if not self.pptx_path.get():
            self.status_bar.config(text="错误：请先选择文件！", foreground="red")
            return
            
        self.running = True
        self.status_bar.config(text="分析进行中...", foreground="blue")
        
        # 在后台线程运行处理任务
        def task():
            try:
                # ... (保持原有处理逻辑不变)
                file_list = ['KineticStandard.xlsx','SteadyStandard.xlsx' ]
                #file_list = ['20250307-Kinetic Standard Curve-Middle Conc.xlsx', '377-Standard Curve.xlsx']

                #pptx_file = '20220809-curveoutput-XHB-CPD1-1280-only multiple curve.pptx'
                #pptx_file = 'chart.pptx'
                data_frame,df_original=self.extract_chart_data(self.pptx_path.get())

                for file in file_list:
                    data = self.read_column_to_list(file,column_name='Y-axis')
                    data_trim = self.retain_numbers_uniform(data,int(self.sample_count.get()))
                    #print(len(data_trim))
                    #print(data_trim)

                    data_frame[file.split()[0]]= data_trim

                data_frame.to_csv('chart_data.csv', index=False)
                
                df_original.to_csv('chart_original_data.csv', index=False)
                
                data_dict = {} 
                columns_data = data_frame.columns[:-2]

                for no in columns_data:
                    data_list =  [float(x) for x in data_frame[no].tolist()]
                    Kinetic_list = [float(x) for x in data_frame[data_frame.columns[-2]].tolist()]
                    Steady_list =  [float(x) for x in data_frame[data_frame.columns[-1]].tolist()]
                    Kinetic_distance=round(data_list[-1]-Kinetic_list[-1],2)
                    Kinetic_MSE =round(self.calculate_mse(Kinetic_list,data_list,Kinetic_distance),2)
                    Steady_distance=round(data_list[-1]-Steady_list[-1],2)
                    Steady_MSE=round(self.calculate_mse(Steady_list,data_list,Steady_distance),2) 
                    #data_dict[no]={'Kinetic_distance':Kinetic_distance,'Kinetic_score':Kinetic_MSE,'Steady_distance':Steady_distance,'Steady_score':Steady_MSE}
                    data_dict[no]={'binding_max':data_list[-1],'Kinetic_score':Kinetic_MSE,'Steady_score':Steady_MSE}
                df_result = pd.DataFrame.from_dict(data_dict, orient="index")
                print(df_result)
                df_result.to_csv('chart_result.csv', index=True)

                self.status_bar.config(text="分析完成", foreground="green")
            except Exception as e:
                self.async_redirect.queue.put(f"\n错误：{str(e)}\n")
            finally:
                self.running = False
                self.root.event_generate("<<TaskComplete>>")
        
        threading.Thread(target=task, daemon=True).start()
        
        # 绑定任务完成事件
        self.root.bind("<<TaskComplete>>", self.on_task_complete)

    def on_task_complete(self, event):
        self.status_bar.config(text="分析完成", foreground="green")
        self.root.unbind("<<TaskComplete>>")

    def __del__(self):
        if hasattr(self, 'async_redirect'):
            self.async_redirect.stop()

    def clear_output(self):
        self.output_text.configure(state='normal')
        self.output_text.delete(1.0, tk.END)
        self.output_text.configure(state='disabled')

    def software_discription(self):        
	    messagebox.showinfo('使用要求',"一、输入文件要求\n\
标准曲线配置文件:\n\
需在工作目录下提供以下两个Excel标准文件：\n\
▸ KineticStandard.xlsx（动态标准曲线）\n\
▸ SteadyStandard.xlsx（稳态标准曲线）\n\
配置文件规范：\n\
每个文件需包含**≥默认48个有效数据点**\n\
支持通过修改数据列更新标准曲线基准值\n\
待分析文件\n\
输入文件格式：PPTX格式演示文稿\n\
内容要求：包含待分析的可视化图表\n\
二、数据处理\n\
生成中间数据文件：\n\
▸ chart_original_data.csv（完整数据集存储）\n\
▸ chart_data.csv（记录峰值前默认48个连续数据点）\n\
输出分析结果文件：\n\
▸ chart_result.csv  ")

if __name__ == "__main__":
    root = tk.Tk()
    app = ModernPPTXApp(root)
    root.mainloop()
    sys.stdout = sys.__stdout__